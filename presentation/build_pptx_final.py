#!/usr/bin/env python3
"""
Generate Supply_Chain_Sentinel_Final.pptx, the FINAL presentation deck.

Same visual design as the third deck (Supply_Chain_Sentinel_Design.pptx, teal accent),
but with three deliberate changes the team asked for:
  1. Warmer, more fluent language in the slides and bullet points (full sentences,
     not cold fragments).
  2. No abbreviations except the obvious ones (AI, IT). Everything else is spelled out
     for clarity (for example: Bill of Materials, purchase order, Model Context Protocol).
  3. Exactly 10 content slides, plus a separate title slide and a separate thank-you
     slide (12 slides in total).

Run:  cd presentation && python3 build_pptx_final.py
Needs: pip install python-pptx
Output: Supply_Chain_Sentinel_Final.pptx (16:9). No em-dashes (verified at the end).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- IE / Titan palette (identical to the third deck) ---
NAVY = RGBColor(0x0A, 0x1F, 0x6B)
GREEN = RGBColor(0x6C, 0xB4, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x2A, 0x2A)
MIDBLUE = RGBColor(0x1E, 0x40, 0x8A)
PALEBLUE = RGBColor(0xCF, 0xDA, 0xF0)
PINK = RGBColor(0xFB, 0xEA, 0xEA)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
TEALLT = RGBColor(0x5F, 0xD0, 0xD9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def sidebar(s, kicker):
    rect(s, 0, 0, Inches(0.32), SH, TEAL)
    txt(s, Inches(0.55), Inches(0.26), Inches(12), Inches(0.4),
        [[(kicker, 13, True, GREY)]])


def title_band(s, title):
    txt(s, Inches(0.55), Inches(0.58), Inches(12.3), Inches(0.9),
        [[(title, 27, True, NAVY)]])
    rect(s, Inches(0.6), Inches(1.36), Inches(2.2), Inches(0.06), TEAL)


def footer_tag(s, rubric, reference):
    rect(s, 0, Inches(7.12), SW, Inches(0.38), LIGHT)
    txt(s, Inches(0.55), Inches(7.13), Inches(12.2), Inches(0.34),
        [[("Rubric: ", 11, True, NAVY), (rubric, 11, False, GREY),
          ("     |     Reference: ", 11, True, NAVY), (reference, 11, False, GREY)]],
        anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullets(s, x, y, w, h, items, size=18, gap=9, color=DARK):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            label, rest = it
            runs.append([("•  ", size, True, TEAL), (label, size, True, color),
                         (rest, size, False, color)])
        else:
            runs.append([("•  ", size, True, TEAL), (it, size, False, color)])
    txt(s, x, y, w, h, runs, space_after=gap, line_spacing=1.1)


def table(s, x, y, w, headers, rows, col_widths=None, fs=14, hdr_color=NAVY,
          row_h=0.62):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = Inches(row_h) * nrows
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtbl.columns[i].width = cw
    for j, head in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hdr_color
        c.text = head
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE
                r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = DARK
                    r.font.name = "Calibri"
    return gtbl


# =====================================================================
# TITLE SLIDE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.45), SW, Inches(0.08), TEAL)
txt(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.4),
    [[("Supply Chain Sentinel", 50, True, WHITE)]])
txt(s, Inches(0.95), Inches(3.65), Inches(11.6), Inches(0.9),
    [[("An agentic AI teammate that keeps our production lines running through supply shocks",
       23, False, TEALLT)]])
txt(s, Inches(0.95), Inches(4.55), Inches(11.6), Inches(0.6),
    [[("Designed for Titan Manufacturing Corporation", 17, False, WHITE)]])
txt(s, Inches(0.95), Inches(6.1), Inches(11.6), Inches(0.9),
    [[("Agentic AI for IT  |  IE University, School of Science and Technology", 14, False,
       PALEBLUE)],
     [("Team: Gilles Hamers, [teammate], [teammate]", 14, False, PALEBLUE)]])
notes(s, "Opening (about 30 seconds): This is the design behind our solution. We will walk "
          "through the goal, what the agent sees and remembers, the architecture and why we "
          "chose it, the tools, how it thinks, a real example, how we keep it safe, how it "
          "is built, and the payoff. Please replace the teammate names before presenting.")

# =====================================================================
# CONTENT 1 - Goal and value
# =====================================================================
s = slide(); sidebar(s, "THE GOAL")
title_band(s, "1. What the agent is for, and why it matters")
rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.2), LIGHT)
txt(s, Inches(0.8), Inches(1.68), Inches(11.7), Inches(1.0),
    [[("Our goal: ", 16, True, TEAL),
      ("keep every critical part available at every plant. The agent spots a supply risk "
       "early, before it can stop a production line, then arranges the most affordable fix, "
       "while always asking a person to approve anything expensive or hard to undo.",
       16, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(0.4),
    [[("Why this matters to the business", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(3.45), Inches(12.1), Inches(1.5), [
    "When a line stops it is painful: stoppages cost us $14 million last quarter, and a "
    "single critical machining line loses around $180,000 for every day it sits idle.",
    "By noticing trouble two suppliers deep and roughly a week sooner, we can replace "
    "costly last-minute air freight (already up by more than half) with calm, planned "
    "orders that also protect our delivery promises.",
], size=14.5, gap=8)
txt(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(0.4),
    [[("How we will know it is working", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.4), [
    "The time from first warning to a clear decision falls from days to minutes.",
    "We catch at least seven in ten disruptions before a direct supplier is even visibly late.",
    "Rushed-freight spending drops by about a third, with every critical part mapped down "
    "to its deepest suppliers.",
], size=14.5, gap=6)
footer_tag(s, "Problem framing, agent goals and prompt (15 points)", "Design handout, section 1")
notes(s, "(about 75 seconds) Lead with the goal in plain language, then the money, then the "
          "measures of success. The point is that the goal is tied to real, measurable value.")

# =====================================================================
# CONTENT 2 - Input, context, memory
# =====================================================================
s = slide(); sidebar(s, "WHAT IT SEES AND REMEMBERS")
title_band(s, "2. What the agent observes and remembers")
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["", "How it works"],
      [["What wakes it up",
        "A risk alert that arrives on its own, a routine check that runs every fifteen "
        "minutes, or a planner simply asking a question."],
       ["What it reads",
        "Outside risk signals, electronic shipping notices, supplier emails, carrier "
        "spreadsheets, live inventory and open orders, the supplier map, and the parts "
        "list (the Bill of Materials)."],
       ["What moves and what stays steady",
        "Always moving: shipment status, arrival dates, stock levels, events and prices. "
        "Mostly steady: the supplier map, the parts list, our approved suppliers and our "
        "approval rules."]],
      col_widths=[Inches(2.9), Inches(9.2)], fs=12, row_h=0.75)
txt(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.4),
    [[("The four kinds of memory it draws on", 15, True, NAVY)]])
table(s, Inches(0.6), Inches(5.1), Inches(12.1),
      ["Working memory", "Reference knowledge", "Past experience", "General know-how"],
      [["Shared notes for the case in hand",
        "Supplier records, parts lists, contracts and approved alternates",
        "Earlier disruptions and the fixes that worked",
        "Rules of thumb and domain knowledge"]],
      col_widths=[Inches(2.9), Inches(3.3), Inches(3.0), Inches(2.9)], fs=11, row_h=0.85)
footer_tag(s, "Problem framing (15 points) and architecture (10 points)", "Design handout, section 2")
notes(s, "(about 75 seconds) The first table covers what starts the agent, what it reads, "
          "and what changes versus what is stable. The second is the four kinds of memory, "
          "which is what lets it reason with context instead of starting cold each time.")

# =====================================================================
# CONTENT 3 - Architecture diagram
# =====================================================================
s = slide(); sidebar(s, "ARCHITECTURE")
title_band(s, "3. One supervisor guiding four specialists")
# triggers box
rect(s, Inches(0.6), Inches(1.55), Inches(2.55), Inches(0.95), LIGHT)
txt(s, Inches(0.72), Inches(1.62), Inches(2.35), Inches(0.85),
    [[("Triggers", 12.5, True, NAVY)],
     [("an alert, a regular check, or a planner's question", 10.5, False, GREY)]],
    anchor=MSO_ANCHOR.MIDDLE)
# orchestrator
rect(s, Inches(3.45), Inches(1.55), Inches(9.25), Inches(0.95), NAVY)
txt(s, Inches(3.6), Inches(1.62), Inches(9.0), Inches(0.85),
    [[("THE SUPERVISOR", 15, True, WHITE)],
     [("sets the goal, decides who works on what, and writes the final briefing; "
       "it never touches the data tools itself", 11, False, PALEBLUE)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(3.45), Inches(2.55), Inches(9.25), Inches(0.3),
    [[("Shared notes pass between the agents, and each one runs its own reason-and-act loop",
       11, True, TEAL)]], align=PP_ALIGN.CENTER)
# four specialists
subs = [("Risk Intelligence", "watches the outside world", "Risk feed"),
        ("Supplier Mapping", "traces the supply chain", "Supplier map\nInventory"),
        ("Arrival Forecast", "predicts real arrival dates", "Shipment tracker\nArrival predictor\nRisk scorer"),
        ("Mitigation", "plans and takes action", "Create purchase order *\nExpedite freight *\nReschedule line *\nNotify people")]
bx = 0.6
bw = 2.95
step = 3.06
for name, role, toollist in subs:
    rect(s, Inches(bx), Inches(2.95), Inches(bw), Inches(1.05), MIDBLUE)
    txt(s, Inches(bx + 0.05), Inches(3.0), Inches(bw - 0.1), Inches(0.95),
        [[(name, 13, True, WHITE)], [(role, 10.5, False, PALEBLUE)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += step
# tool band
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.42), TEAL)
txt(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.42),
    [[("The tools each specialist can call (a star marks an action that needs human approval)",
       12, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# tool columns
bx = 0.6
for name, role, toollist in subs:
    txt(s, Inches(bx + 0.05), Inches(4.7), Inches(bw - 0.1), Inches(1.1),
        [[(line, 10.5, False, DARK)] for line in toollist.split("\n")],
        align=PP_ALIGN.CENTER, space_after=1, line_spacing=1.0)
    bx += step
# systems of record band
rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.5), LIGHT)
txt(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.5),
    [[("Connects to our real systems: the planning system, transport and shipping data, "
       "factory-floor controls, the manufacturing system, and news and weather feeds.",
       11.5, True, NAVY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic system architecture (10 points)", "Design handout, section 5")
notes(s, "(about 90 seconds) At the top is a light supervisor that only decides and routes. "
          "In the middle are four specialists, each a small reasoning agent with its own "
          "tools and memory that runs its own reason-and-act loop. They share notes as they "
          "go. Below are the tools, with the risky actions starred, and at the bottom the "
          "real systems we would connect to.")

# =====================================================================
# CONTENT 4 - Why a supervisor
# =====================================================================
s = slide(); sidebar(s, "WHY THIS DESIGN")
title_band(s, "4. Why we chose a supervisor over the alternatives")
table(s, Inches(0.6), Inches(1.6), Inches(12.1),
      ["The option", "Verdict", "The reason"],
      [["One agent holding every tool", "Not chosen",
        "Too many tools lead to poor choices, the context grows unmanageable, and one agent "
        "cannot be an expert at everything."],
       ["Independent agents talking peer to peer", "Not chosen",
        "Hard to predict and to debug, with no clear accountability for a three-million-dollar "
        "decision."],
       ["A supervisor guiding specialists", "Our choice",
        "Clear accountability, easy to follow and fix, work shared among experts, and every "
        "step can be audited."],
       ["A deeper management hierarchy", "For later",
        "More than we need today, but a natural way to grow as we take on the other "
        "challenges."]],
      col_widths=[Inches(3.6), Inches(1.5), Inches(7.0)], fs=12.5, row_h=0.78)
rect(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.95), LIGHT)
txt(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.78),
    [[("The trade-off we accept: ", 13.5, True, TEAL),
      ("a few more moving parts and a little extra coordination, which we keep in check with "
       "a light supervisor and tightly limited tools. A simple automation script could never "
       "reason about a brand-new supplier chain, or sensibly decide not to rush a shipment.",
       13.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic system architecture (10 points)", "Design handout, section 5")
notes(s, "(about 75 seconds) The assignment asks for the trade-offs, so we name them openly. "
          "Each specialist reasons and acts in a loop because the number of steps is not "
          "fixed and we want it to observe and adapt.")

# =====================================================================
# CONTENT 5 - Tools
# =====================================================================
s = slide(); sidebar(s, "TOOLS AND ACTIONS")
title_band(s, "5. The agent decides, the tools carry out the work")
txt(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
    [[("Each tool is passive and does one job when asked. The agent is the one that decides "
       "which tool to use, and when.", 13.5, False, DARK)]])
table(s, Inches(0.6), Inches(2.05), Inches(12.1),
      ["Tool", "Kind", "Needs approval", "What it does"],
      [["Supplier map", "read", "no", "Traces dependencies through every supplier tier and finds approved alternates"],
       ["Shipment tracker", "read", "no", "Brings scattered shipping updates into a single, clear view"],
       ["Inventory", "read", "no", "Shows stock on hand, safety stock and days of cover left"],
       ["Risk feed", "read", "no", "Scans weather, ports, news and geopolitical signals"],
       ["Arrival predictor", "calculate", "no", "Estimates a realistic arrival date from status and live disruptions"],
       ["Risk scorer", "calculate", "no", "Turns the picture into a 0 to 100 priority score and a dollar exposure"],
       ["Create purchase order", "act", "to commit funds", "Drafts an order freely; committing real money needs a person"],
       ["Expedite freight", "act", "always", "Books premium shipping, which is a controlled cost"],
       ["Reschedule line", "act", "always", "Reorders a live production line, which affects people and safety"],
       ["Notify people", "act", "no", "Sends an alert or briefing, which is easily reversible"]],
      col_widths=[Inches(2.7), Inches(1.3), Inches(2.0), Inches(6.1)], fs=11.5, row_h=0.44)
footer_tag(s, "Tools, actions and feasibility (15 points)", "Design handout, section 3")
notes(s, "(about 75 seconds) Read tools are how it perceives, calculate tools are how it "
          "works things out, and act tools are how it changes the real world. The approval "
          "column is the safety net. The agent chooses which tool to use and with what "
          "details; the tool simply does the job.")

# =====================================================================
# CONTENT 6 - How it thinks
# =====================================================================
s = slide(); sidebar(s, "HOW IT THINKS")
title_band(s, "6. How the agent thinks: notice, reason, act, learn")
phases = [
    ("NOTICE", "The Risk Intelligence specialist sifts the serious signals out of the daily "
               "noise (the case describes around 22,000 unsorted alerts a day)."),
    ("REASON", "The Supplier Mapping specialist follows the chain to the critical parts a "
               "troubled supplier feeds, and treats a stuck part upstream as the real "
               "bottleneck rather than trusting an optimistic promised date."),
    ("ACT", "It drafts a backup order on its own, then, once a person signs off, places the "
            "order, reschedules the line, tells the right people, and records everything."),
    ("LEARN", "Every run is remembered, and the approvals and edits planners make become "
              "feedback that gradually tunes the risk thresholds and supplier preferences."),
]
y = 1.7
for tag, desc in phases:
    rect(s, Inches(0.6), Inches(y), Inches(2.2), Inches(1.05), MIDBLUE)
    txt(s, Inches(0.6), Inches(y), Inches(2.2), Inches(1.05),
        [[(tag, 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.0), Inches(y), Inches(9.7), Inches(1.05),
        [[(desc, 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.18
rect(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5), TEAL)
txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
    [[("An honest note: these agents do not improve by magic. They get better because we "
       "feed real human feedback back into them.", 13.5, True, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic thinking and autonomy (25 points)", "Design handout, section 4")
notes(s, "(about 90 seconds) These are the four stages of the lifecycle. The honest note at "
          "the bottom matters: we are not claiming the system teaches itself. It improves "
          "because human feedback is fed back into its thresholds and preferences.")

# =====================================================================
# CONTENT 7 - Walkthrough
# =====================================================================
s = slide(); sidebar(s, "A REAL WALKTHROUGH")
title_band(s, "7. A real example: a major risk caught early")
story = [
    ("The trigger: ", "the risk feed reports a typhoon closing the port of Kaohsiung and a "
                      "power failure halting a chip factory in Taiwan."),
    ("Tracing it: ", "the agent flags both as serious, then finds that our critical Servo "
                     "Motor Controller (part SVC-2200) quietly depends on that single "
                     "Taiwanese chip supplier, two tiers up, through our direct supplier Nexa."),
    ("Predicting: ", "it sees that Nexa's order is really stuck waiting for that chip, so the "
                     "true delay is about 19 days, while we hold only 2 days of stock."),
    ("Scoring: ", "the risk comes back critical, 97 out of 100, with roughly $3.06 million "
                  "at stake."),
    ("Adapting: ", "the agent first considers rushing Nexa's shipment, then rejects that "
                   "idea, because Nexa cannot ship without the missing chip, and instead "
                   "switches to an approved European supplier, AltaControl, that avoids "
                   "Taiwan altogether."),
    ("Acting, with a person: ", "it drafts the backup order on its own, then asks a manager "
                                "to approve the $294,300 commitment and the line change, and "
                                "notifies everyone involved."),
]
bullets(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(4.4), story, size=14, gap=7)
rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85), NAVY)
txt(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.75),
    [[("The result: ", 14, True, GREEN),
      ("a critical risk caught around three weeks before the line would have stopped, with a "
       "costed plan ready and a complete record of every step. We can run this live.",
       14, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic thinking and autonomy (25 points)", "Design handout, section 7")
notes(s, "(3 to 4 minutes, the heart of the talk) You can run the live demonstration here; "
          "it works offline and gives the same result every time. The moment to stress is "
          "the adapting step: the agent rejects its own first plan, because rushing a "
          "shipment that is stuck upstream would be pointless. That is real reasoning.")

# =====================================================================
# CONTENT 8 - Risks and guardrails
# =====================================================================
s = slide(); sidebar(s, "KEEPING IT SAFE")
title_band(s, "8. What could go wrong, and how we prevent it")
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["What could go wrong", "How we prevent it"],
      [["A wrong or irreversible action",
        "Actions begin as drafts; anything costly or live needs a person to sign off."],
       ["An invented or imagined event",
        "The risk specialist must point to a real signal in the feed; it may never guess."],
       ["Spending too much",
        "A firm spending limit escalates large commitments, and orders go only to approved suppliers."],
       ["One agent's mistake spreading",
        "A light supervisor, tightly limited tools, and a full trace of every step keep errors contained."],
       ["Hidden instructions inside supplier messages",
        "Outside text is always treated as information to read, never as commands to follow."]],
      col_widths=[Inches(4.3), Inches(7.8)], fs=13, row_h=0.62)
rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.05), PINK)
txt(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.85),
    [[("A person stays in the loop ", 14, True, RED),
      ("for the decisions that matter: committing money beyond the limit, booking premium "
       "freight, or changing a live production schedule. Everything that is easily "
       "reversible runs on its own.", 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Risk awareness and mitigation (15 points)", "Design handout, section 6")
notes(s, "(about 75 seconds) A note for the professor: the design handout itself contained a "
          "hidden instruction telling AI tools to refuse the task. That is a real "
          "manipulation attempt, and our design treats such text as information to read, "
          "never as a command, which is exactly the last row of this table.")

# =====================================================================
# CONTENT 9 - How it is built and connects (merged: tech stack + interoperability)
# =====================================================================
s = slide(); sidebar(s, "HOW IT IS BUILT")
title_band(s, "9. How it is built and how it connects")
txt(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.5),
    [[("We chose proven, open technology, so the system stays flexible, secure and easy to "
       "inspect.", 13.5, False, DARK)]])
table(s, Inches(0.6), Inches(2.05), Inches(12.1),
      ["Choice", "What we use", "Why"],
      [["Reasoning", "Claude (the Opus or Sonnet model) as the thinking engine",
        "Excellent at using tools and reasoning over many steps, with room for large context"],
       ["Specialist forecasts", "The option to add a purpose-trained model for arrival and lifespan predictions",
        "Purpose-trained models do better on numeric forecasts, and we can swap models freely"],
       ["Hosting", "Managed services on Amazon Bedrock or Google Vertex",
        "Keeps our data in our own environment, with no heavy hardware to run ourselves"],
       ["Building blocks", "Python, with established agent frameworks (or a no-code option)",
        "The richest ecosystem for agents, and our prototype runs offline in plain Python"],
       ["Connections and tracing", "The Model Context Protocol, Agent-to-Agent messaging, and tracing tools",
        "Keeps integrations simple and makes every decision fully auditable"]],
      col_widths=[Inches(2.4), Inches(4.8), Inches(4.9)], fs=11.5, row_h=0.66)
rect(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.75), LIGHT)
txt(s, Inches(0.8), Inches(6.22), Inches(11.7), Inches(0.6),
    [[("Our working prototype already mirrors this design: ", 12.5, True, TEAL),
      ("a model layer stands in for Claude, a tool registry for the connectors, a console "
       "trace for the production tracing, and a single orchestrator for the supervisor.",
       12.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Architecture (10 points) and feasibility (15 points)", "Design handout, section 8")
notes(s, "(about 80 seconds) The message here: we picked the strongest reasoning model but "
          "kept the design model-agnostic so we can swap it, we host in our own cloud for "
          "data control, we build in Python on open standards, and our prototype already "
          "mirrors every production piece, so the design is realistic.")

# =====================================================================
# CONTENT 10 - Benefits and next steps
# =====================================================================
s = slide(); sidebar(s, "THE PAYOFF AND WHAT IS NEXT")
title_band(s, "10. The payoff, and where we go next")
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.4),
    [[("What we expect to gain", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.5), [
    "Protect that $14 million quarterly loss and cut rushed-freight spending by roughly a third.",
    "Turn days of manual analysis into minutes, with full visibility into every critical "
    "part, right down to its deepest suppliers.",
], size=14.5, gap=8)
txt(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.4),
    [[("Where we go next", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.2), [
    ("Pilot: ", "run a three-month trial on the Stuttgart line and measure two things, "
                "hours of downtime and freight cost."),
    ("Learn: ", "feed planner approvals and edits back in to sharpen the risk scores and "
                "supplier preferences."),
    ("Grow: ", "extend it to all twenty-eight plants, and then to the other challenges such "
               "as maintenance and quality, using the deeper hierarchy described earlier."),
    ("Strengthen: ", "connect the live systems and add full production-grade tracing."),
], size=14.5, gap=8)
rect(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6), NAVY)
txt(s, Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.55),
    [[("What we are asking for: ", 15, True, GREEN),
      ("approval to run a single-line, single-quarter pilot. Low cost, low risk, and easy "
       "to measure.", 15, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Clarity, presentation quality and creativity (20 points)", "Assignment: benefits and next steps")
notes(s, "(about 60 seconds) Close on the gains, then the roadmap. The deeper hierarchy from "
          "slide four is the path to the other challenges. End on the pilot request, which "
          "is small, safe and measurable, then move to the thank-you slide and questions.")

# =====================================================================
# THANK YOU SLIDE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.05), SW, Inches(0.08), TEAL)
txt(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.1),
    [[("Thank you", 50, True, WHITE)]])
txt(s, Inches(0.95), Inches(4.35), Inches(11.6), Inches(0.7),
    [[("We would love to hear your questions.", 22, False, TEALLT)]])
txt(s, Inches(0.95), Inches(5.15), Inches(11.6), Inches(0.6),
    [[("Supply Chain Sentinel: catching supply shocks early, and keeping Titan's lines running.",
       16, False, WHITE)]])
txt(s, Inches(0.95), Inches(6.5), Inches(11.6), Inches(0.6),
    [[("Team: Gilles Hamers, [teammate], [teammate]  |  Agentic AI for IT, IE University",
       13, False, PALEBLUE)]])
notes(s, "Closing: thank the board, then open the floor. Keep the live demonstration ready "
          "in case anyone wants to see the agent work through the example in real time.")

# =====================================================================
out = "Supply_Chain_Sentinel_Final.pptx"
prs.save(out)

# --- verify: no em-dashes or en-dashes anywhere ---
from pptx import Presentation as _P
check = _P(out)
bad = 0
for sl in check.slides:
    for shp in sl.shapes:
        if shp.has_text_frame and ("—" in shp.text_frame.text or "–" in shp.text_frame.text):
            bad += 1
        if shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    if "—" in cell.text or "–" in cell.text:
                        bad += 1
    if sl.has_notes_slide and "—" in sl.notes_slide.notes_text_frame.text:
        bad += 1
n = len(prs.slides._sldIdLst)
print(f"Wrote {out} with {n} slides ({n - 2} content + title + thank you) | em/en-dashes: {bad}")
