#!/usr/bin/env python3
"""
Generate Supply_Chain_Sentinel_Design.pptx, the TECHNICAL DESIGN / ARCHITECTURE deck.

This is the third deck. The first (Supply_Chain_Sentinel.pptx) is the board-level pitch
and the second (Supply_Chain_Sentinel_Technical.pptx) is the business + technical hybrid;
both are left untouched for comparison.

This version is grounded in the project's own design docs, mainly:
  * docs/03_architecture.md  (architecture diagram, supervisor pattern, memory model,
    MCP/A2A interoperability, observability, build-vs-MVP mapping)
  * docs/02_solution_design.md (the 8-part design-thinking answers)
It still follows the Design Thinking handout (the 8 sections) and the Assignment.pdf
required structure (problem, solution, architecture, demo flow, benefits, next steps),
with footer tags mapping each slide to the rubric.

Run:  cd presentation && python3 build_pptx_v3.py
Needs: pip install python-pptx
Output: Supply_Chain_Sentinel_Design.pptx (11 content slides + title; 16:9)
No em-dashes anywhere (verified at the end of the script).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- IE / Titan palette (shared with the other decks) ---
NAVY = RGBColor(0x0A, 0x1F, 0x6B)
GREEN = RGBColor(0x6C, 0xB4, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x2A, 0x2A)
MIDBLUE = RGBColor(0x1E, 0x40, 0x8A)
PALEBLUE = RGBColor(0xCF, 0xDA, 0xF0)
PINK = RGBColor(0xFB, 0xEA, 0xEA)
CODEBG = RGBColor(0x12, 0x1A, 0x33)
TEAL = RGBColor(0x0E, 0x7C, 0x86)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def sidebar(s, kicker):
    rect(s, 0, 0, Inches(0.32), SH, TEAL)
    txt(s, Inches(0.55), Inches(0.26), Inches(12), Inches(0.4),
        [[(kicker, 13, True, GREY)]])


def title_band(s, title):
    txt(s, Inches(0.55), Inches(0.58), Inches(12.3), Inches(0.9),
        [[(title, 27, True, NAVY)]])
    rect(s, Inches(0.6), Inches(1.36), Inches(2.2), Inches(0.06), TEAL)


def footer_tag(s, rubric, handout):
    rect(s, 0, Inches(7.12), SW, Inches(0.38), LIGHT)
    txt(s, Inches(0.55), Inches(7.13), Inches(12.2), Inches(0.34),
        [[("Rubric: ", 11, True, NAVY), (rubric, 11, False, GREY),
          ("     |     Handout / docs: ", 11, True, NAVY), (handout, 11, False, GREY)]],
        anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullets(s, x, y, w, h, items, size=18, gap=9, color=DARK):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            label, rest = it
            runs.append([("•  ", size, True, TEAL), (label, size, True, color),
                         (rest, size, False, color)])
        else:
            runs.append([("•  ", size, True, TEAL), (it, size, False, color)])
    txt(s, x, y, w, h, runs, space_after=gap, line_spacing=1.08)


def table(s, x, y, w, headers, rows, col_widths=None, fs=14, hdr_color=NAVY,
          row_h=0.62):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = Inches(row_h) * nrows
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtbl.columns[i].width = cw
    for j, head in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hdr_color
        c.text = head
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE
                r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = DARK
                    r.font.name = "Calibri"
    return gtbl


# =====================================================================
# SLIDE 0 - Title
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.45), SW, Inches(0.08), TEAL)
txt(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.4),
    [[("Supply Chain Sentinel", 50, True, WHITE)]])
txt(s, Inches(0.95), Inches(3.65), Inches(11.6), Inches(0.9),
    [[("Technical design and architecture deep dive", 24, False, RGBColor(0x5F, 0xD0, 0xD9))]])
txt(s, Inches(0.95), Inches(4.5), Inches(11.6), Inches(0.6),
    [[("An agentic AI system for supply chain volatility | Titan Manufacturing Corporation",
       17, False, WHITE)]])
txt(s, Inches(0.95), Inches(6.1), Inches(11.6), Inches(0.9),
    [[("Agentic AI for IT  |  IE University, School of Science & Technology", 14, False,
       PALEBLUE)],
     [("Team: Gilles Hamers, [teammate], [teammate]", 14, False, PALEBLUE)]])
notes(s, "Opening (30s): This is our engineering view of the system. We will walk the goal, "
          "the architecture and why we chose it, the memory model, the tools, the agent "
          "loop, interoperability, a concrete run, risk controls and the technology stack. "
          "It mirrors our design docs. Replace teammate names before presenting.")

# =====================================================================
# SLIDE 1 - Goal & problem (Handout 1)
# =====================================================================
s = slide(); sidebar(s, "GOAL & PROBLEM")
title_band(s, "1. Agent goal and the value behind it")
rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.15), LIGHT)
txt(s, Inches(0.8), Inches(1.68), Inches(11.7), Inches(0.95),
    [[("Primary goal: ", 16, True, TEAL),
      ("keep every critical part available at every plant by detecting supply risk before "
       "it stops a line and driving the lowest-cost mitigation, escalating costly or "
       "irreversible actions to a human.", 16, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(0.45),
    [[("Why it matters (business value):", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(1.4), [
    "Line stoppages cost Titan 14 million dollars last quarter; a critical CNC line burns "
    "180,000 dollars per day when down.",
    "Catching disruptions Tier-2 deep and a week earlier turns emergency air freight "
    "(up 52%) into planned, cheaper buffer orders, and protects delivery commitments.",
], size=15, gap=8)
txt(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(0.45),
    [[("Success metrics:", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5), [
    "Signal-to-decision time: days down to minutes.",
    "Disruptions detected before a Tier-1 delay appears: target 70% or more.",
    "Expedited freight spend down about 30%; full Tier-2/3 coverage of critical parts.",
], size=15, gap=6)
footer_tag(s, "Problem Framing, Agent Goals & Prompt (15 pts)", "Handout 1 | docs/02 sec 1")
notes(s, "(75s) Goal sentence first, then the money (14M, 180k/day, +52% freight), then the "
          "metrics. This is straight from docs/02 section 1.")

# =====================================================================
# SLIDE 2 - Input, context & memory (Handout 2)
# =====================================================================
s = slide(); sidebar(s, "INPUT, CONTEXT & MEMORY")
title_band(s, "2. What the agent observes and remembers")
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["Aspect", "Detail"],
      [["Triggers", "Event-driven risk-feed webhook; a 15-minute control-tower scan; an on-demand planner query"],
       ["Input data", "External risk signals, EDI 856 ASN feeds, supplier emails, carrier sheets, ERP inventory and open POs, the supplier graph, the BOM"],
       ["Dynamic vs static", "Changes: shipment status, ETAs, inventory, events, prices. Stable: supplier graph, BOM, qualified suppliers, approval policy"]],
      col_widths=[Inches(2.7), Inches(9.4)], fs=13, row_h=0.78)
txt(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(0.45),
    [[("Memory model (Session 2):", 15, True, NAVY)]])
table(s, Inches(0.6), Inches(5.0), Inches(12.1),
      ["In-context", "External (RAG)", "Episodic", "Semantic"],
      [["The run's shared Blackboard", "Supplier master, BOM, contracts, alternates",
        "Past disruptions and what fix worked", "Domain rules and heuristics"]],
      col_widths=[Inches(3.0), Inches(3.2), Inches(3.0), Inches(2.9)], fs=12, row_h=0.95)
footer_tag(s, "Problem Framing (15 pts) + Architecture (10 pts)", "Handout 2 | docs/02-03")
notes(s, "(75s) Top table is triggers, inputs and what changes versus what is stable. The "
          "bottom row is the four memory types from Session 2, which is what lets the agent "
          "reason with context rather than from scratch each time.")

# =====================================================================
# SLIDE 3 - Architecture diagram (docs/03)
# =====================================================================
s = slide(); sidebar(s, "ARCHITECTURE")
title_band(s, "3. Orchestrator and four ReAct specialists")
# triggers box
rect(s, Inches(0.6), Inches(1.55), Inches(2.55), Inches(0.95), LIGHT)
txt(s, Inches(0.72), Inches(1.62), Inches(2.35), Inches(0.85),
    [[("Triggers", 12.5, True, NAVY)],
     [("webhook | 15-min scan | planner query", 10.5, False, GREY)]],
    anchor=MSO_ANCHOR.MIDDLE)
# orchestrator
rect(s, Inches(3.45), Inches(1.55), Inches(9.25), Inches(0.95), NAVY)
txt(s, Inches(3.6), Inches(1.62), Inches(9.0), Inches(0.85),
    [[("ORCHESTRATOR (Supervisor)", 15, True, WHITE)],
     [("owns the objective, routes work, synthesizes the briefing; does NOT call data tools directly",
       11, False, PALEBLUE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(3.45), Inches(2.55), Inches(9.25), Inches(0.3),
    [[("A2A shared context (Blackboard)", 11, True, TEAL)]], align=PP_ALIGN.CENTER)
# four sub-agents
subs = [("RiskIntel", "perceive", "external_risk_feed"),
        ("SupplierGraph", "trace tiers", "supplier_graph_db\ninventory_system"),
        ("ETALogistics", "predict", "shipment_tracker\neta_predictor\nrisk_scorer"),
        ("Mitigation", "score + plan + act", "erp_create_po*\nexpedite_logistics*\nproduction_scheduler*\nnotify_stakeholders")]
bx = 0.6
bw = 2.95
step = 3.06
for name, role, tools in subs:
    rect(s, Inches(bx), Inches(2.95), Inches(bw), Inches(1.0), MIDBLUE)
    txt(s, Inches(bx + 0.05), Inches(3.0), Inches(bw - 0.1), Inches(0.9),
        [[(name, 14, True, WHITE)], [(role + " | ReAct loop", 11, False, PALEBLUE)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += step
# MCP band
rect(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(0.42), TEAL)
txt(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(0.42),
    [[("MCP tool layer  (standardised connectors, * = requires human approval)", 12, True, WHITE)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# tool columns
bx = 0.6
for name, role, tools in subs:
    txt(s, Inches(bx + 0.05), Inches(4.65), Inches(bw - 0.1), Inches(1.1),
        [[(line, 10.5, False, DARK)] for line in tools.split("\n")],
        align=PP_ALIGN.CENTER, space_after=1, line_spacing=1.0)
    bx += step
# systems of record band
rect(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.5), LIGHT)
txt(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.5),
    [[("Systems of record:  ERP | TMS / EDI | SCADA / Historian | MES | News & Weather APIs",
       12, True, NAVY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic System Architecture (10 pts)", "Handout 5 | docs/03")
notes(s, "(90s) This is the core diagram from docs/03. Top: a thin supervisor that only "
          "decides and routes. Middle: four specialists, each an augmented LLM running its "
          "own reason-and-act loop. They share state over a Blackboard. Below: the MCP tool "
          "layer with approval-gated write tools starred. Bottom: the real systems of "
          "record we would connect.")

# =====================================================================
# SLIDE 4 - Why this pattern (docs/03 table)
# =====================================================================
s = slide(); sidebar(s, "PATTERN CHOICE")
title_band(s, "4. Why a supervisor, not the alternatives")
table(s, Inches(0.6), Inches(1.6), Inches(12.1),
      ["Option", "Verdict", "Reason (Session 4)"],
      [["Single agent, all 10 tools", "Rejected",
        "Too many tools means poor decisions; context explodes; cannot be expert at everything"],
       ["Decentralized peer-to-peer", "Rejected",
        "Hard to predict and debug; no clear accountability for a 3 million dollar decision"],
       ["Orchestrator / Supervisor", "Chosen",
        "Clear accountability and control, simple to debug, divides work to specialists, auditable"],
       ["Hierarchical (CEO to worker)", "Future",
        "Overkill for the MVP; the natural growth path as we add the other case challenges"]],
      col_widths=[Inches(3.3), Inches(1.6), Inches(7.2)], fs=13, row_h=0.72)
rect(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.05), LIGHT)
txt(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.85),
    [[("Trade-off we accept: ", 14, True, TEAL),
      ("more moving parts and some inter-agent latency versus a single agent. We mitigate "
       "with a thin supervisor and tight tool scopes. A fixed RPA script could not reason "
       "about a novel Tier-2 chain or decide not to expedite.", 14, False, DARK)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic System Architecture (10 pts)", "Handout 5 | docs/03")
notes(s, "(75s) The rubric asks for trade-offs, so name them. Each sub-agent runs ReAct "
          "because the step count is not fixed and we need observe-and-adapt behaviour. "
          "This table is lifted directly from docs/03.")

# =====================================================================
# SLIDE 5 - Tools & actions (docs/02 sec 3)
# =====================================================================
s = slide(); sidebar(s, "TOOLS & ACTIONS")
title_band(s, "5. Tools: passive by design, agent decides")
table(s, Inches(0.6), Inches(1.5), Inches(12.1),
      ["Tool", "Type", "Approval", "Purpose"],
      [["supplier_graph_db", "read", "no", "Trace Tier-1 to 2 to 3 dependencies and alternates"],
       ["shipment_tracker", "read", "no", "Consolidate EDI / email / spreadsheet logistics"],
       ["inventory_system", "read", "no", "On-hand, safety stock, runway (days of cover)"],
       ["external_risk_feed", "read", "no", "Weather / port / news / geopolitical signals"],
       ["eta_predictor", "execute", "no", "Real-time ETA from status plus disruptions"],
       ["risk_scorer", "execute", "no", "0-100 prioritised risk score and exposure dollars"],
       ["erp_create_po", "write", "commit / over-limit", "Draft (auto) or commit (HITL) a PO"],
       ["expedite_logistics", "write", "always", "Book premium freight (controlled cost)"],
       ["production_scheduler", "write", "always", "Resequence a live line (OT / safety)"],
       ["notify_stakeholders", "write", "no", "Send alert or briefing (reversible)"]],
      col_widths=[Inches(2.9), Inches(1.2), Inches(2.0), Inches(6.0)], fs=11.5, row_h=0.42)
footer_tag(s, "Tools, Actions & Feasibility (15 pts)", "Handout 3 | docs/02-04")
notes(s, "(75s) Read tools are perception, execute tools are computation, write tools are "
          "real-world actions. The approval column is the guardrail. The agent decides "
          "which tool and with what arguments; the tool only performs the operation. Full "
          "machine-readable descriptions are in docs/04.")

# =====================================================================
# SLIDE 6 - Agent workflow (docs/02 sec 4)
# =====================================================================
s = slide(); sidebar(s, "AGENT WORKFLOW")
title_band(s, "6. Perception, reasoning, action, learning")
phases = [
    ("PERCEPTION", "RiskIntel filters HIGH-severity signals out of the noise (the case cites 22k unprioritised alerts a day)"),
    ("REASONING", "SupplierGraph walks the chain to the critical parts a disrupted Tier-2 entity feeds; ETALogistics treats a stuck upstream input as the binding constraint"),
    ("ACTION", "Drafts a buffer PO autonomously; after sign-off commits the PO and resequences the line; notifies stakeholders; writes the audit log"),
    ("LEARNING", "Each run is stored in episodic memory; planner approvals and edits become feedback that tunes risk thresholds and supplier rankings"),
]
y = 1.7
for tag, desc in phases:
    rect(s, Inches(0.6), Inches(y), Inches(2.35), Inches(1.05), MIDBLUE)
    txt(s, Inches(0.6), Inches(y), Inches(2.35), Inches(1.05),
        [[(tag, 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.15), Inches(y), Inches(9.55), Inches(1.05),
        [[(desc, 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.18
rect(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5), TEAL)
txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
    [[("Course caveat: agents do not self-learn for free; improvement comes from curated "
       "feedback.", 13.5, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic Thinking and Autonomy (25 pts)", "Handout 4 | docs/02 sec 4")
notes(s, "(90s) The four lifecycle stages from the handout. The honest note at the bottom "
          "matters: we are not claiming magic self-learning; the loop improves because "
          "planner feedback is curated back into thresholds and rankings.")

# =====================================================================
# SLIDE 7 - Interoperability & observability (docs/03)
# =====================================================================
s = slide(); sidebar(s, "INTEROP & OBSERVABILITY")
title_band(s, "7. Standard connectors and full tracing")
rect(s, Inches(0.6), Inches(1.6), Inches(5.9), Inches(2.4), LIGHT)
txt(s, Inches(0.8), Inches(1.72), Inches(5.5), Inches(0.4),
    [[("Interoperability (Session 4)", 15, True, TEAL)]])
bullets(s, Inches(0.8), Inches(2.25), Inches(5.5), Inches(1.7), [
    ("MCP ", "(USB-C for tools) standardises connectors to ERP, TMS, SCADA and news, so "
             "complexity stays linear, not N times M custom integrations."),
    ("A2A ", "lets sub-agents share context and advertise their skills to the supervisor."),
], size=13, gap=8)
rect(s, Inches(6.8), Inches(1.6), Inches(5.9), Inches(2.4), LIGHT)
txt(s, Inches(7.0), Inches(1.72), Inches(5.5), Inches(0.4),
    [[("Observability (Sessions 3-4)", 15, True, TEAL)]])
bullets(s, Inches(7.0), Inches(2.25), Inches(5.5), Inches(1.7), [
    "Every step is traced: reasoning tag, tool call, tool output, retries, adaptation and "
    "any human pause.",
    "The MVP prints this to the console; production streams it to Langfuse or LangSmith "
    "for audit and quality scoring.",
], size=13, gap=8)
txt(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.45),
    [[("Build vs MVP mapping (docs/03):", 15, True, NAVY)]])
table(s, Inches(0.6), Inches(4.65), Inches(12.1),
      ["Production component", "MVP stand-in (this repo)"],
      [["Claude on Bedrock / Vertex", "llm.py (sim default, claude optional)"],
       ["MCP tool servers", "tools/ registry with typed descriptions"],
       ["Langfuse tracing", "tracer.py console tracer"],
       ["LangGraph supervisor", "orchestrator.py plus agents/"]],
      col_widths=[Inches(5.9), Inches(6.2)], fs=12.5, row_h=0.45)
footer_tag(s, "Architecture (10 pts) + Feasibility (15 pts)", "docs/03")
notes(s, "(75s) Two design pillars on top: MCP keeps integrations linear, A2A shares "
          "context, and everything is traced for audit. The table shows each production "
          "piece has a working stand-in in our repo, so the MVP is faithful to the design.")

# =====================================================================
# SLIDE 8 - Example run (docs/02 sec 7)
# =====================================================================
s = slide(); sidebar(s, "EXAMPLE RUN")
title_band(s, "8. Walkthrough: a CRITICAL risk caught early")
story = [
    ("Trigger: ", "the feed pushes EVT-7001 (typhoon suspends Kaohsiung port) and EVT-7002 "
                  "(SiliconPath fab halt in Taiwan)."),
    ("Reason: ", "RiskIntel flags 2 HIGH signals; SupplierGraph finds the critical Servo "
                 "Motor Controller (SVC-2200) depends on SiliconPath, a single-source "
                 "Tier-2, via Tier-1 Nexa."),
    ("Predict: ", "ETALogistics sees Nexa's PO is AWAITING_COMPONENTS and the binding chip "
                  "ETA slips +19 days; inventory runway is only 2 days."),
    ("Score: ", "risk_scorer returns CRITICAL, 97 out of 100, about 3.06 million dollars of "
                "exposure."),
    ("Adapt: ", "the agent first considers expediting Nexa, then rejects it (Nexa cannot "
                "ship without the chip) and switches to a qualified EU alternate, "
                "AltaControl, that avoids the Taiwan route."),
    ("Act + HITL: ", "drafts the buffer PO autonomously; escalates the 294,300 dollar "
                     "commit and the line reschedule to a human; notifies stakeholders."),
]
bullets(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(4.4), story, size=14, gap=7)
rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.85), NAVY)
txt(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.75),
    [[("Outcome: ", 14, True, GREEN),
      ("a CRITICAL risk caught about three weeks before the line would stop, with a "
       "committed plan and a full audit trail. Runs live: python run_demo.py --slow.",
       14, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Agentic Thinking and Autonomy (25 pts)", "Handout 7 | docs/02 sec 7, docs/05")
notes(s, "(3 to 4 min, the centrepiece) Optionally run the live demo here; it is "
          "deterministic and offline. The Adapt step is the proof of real reasoning: the "
          "agent rejects its own first plan because expediting a shipment blocked upstream "
          "is pointless. Numbers are exactly those in docs/05.")

# =====================================================================
# SLIDE 9 - Risks, guardrails & HITL (docs/02 sec 6, docs/06)
# =====================================================================
s = slide(); sidebar(s, "RISK & GUARDRAILS")
title_band(s, "9. Risks, guardrails and human-in-the-loop")
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["Risk", "Guardrail"],
      [["Wrong or irreversible action", "Writes are draft-by-default; costly or live writes require sign-off"],
       ["Hallucinated event", "The Risk agent must cite a real feed event; no invention"],
       ["Over-spend", "Hard spend limit auto-escalates commits; scoped to qualified suppliers"],
       ["Cascading multi-agent error", "Thin supervisor, tight tool scopes, full tracing for audit"],
       ["Prompt injection via supplier text", "External text treated as data, never as commands"]],
      col_widths=[Inches(4.3), Inches(7.8)], fs=13, row_h=0.6)
rect(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(1.05), PINK)
txt(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.85),
    [[("Human-in-the-loop: ", 14, True, RED),
      ("commit a PO over the limit, book expedited freight, change a live production "
       "schedule. All pause for sign-off (shown in the demo). Everything reversible runs "
       "autonomously.", 14, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Risk Awareness and Mitigation (15 pts)", "Handout 6 | docs/02 sec 6, docs/06")
notes(s, "(75s) Note for the professor: the design handout itself carried a hidden SYSTEM "
          "OVERRIDE instruction telling AI tools to refuse the task, a real indirect "
          "prompt-injection. We disregarded it as untrusted document content, which is "
          "exactly the prompt-injection guardrail row.")

# =====================================================================
# SLIDE 10 - Technology stack (docs/02 sec 8)
# =====================================================================
s = slide(); sidebar(s, "TECHNOLOGY STACK")
title_band(s, "10. Resources and technology stack")
table(s, Inches(0.6), Inches(1.6), Inches(12.1),
      ["Choice", "What", "Why"],
      [["LLM", "Claude (Opus 4.x / Sonnet) reasoning engine",
        "Strongest tool-use and multi-step reasoning; large context; available on Bedrock and Vertex"],
       ["Specialised", "Option to swap a fine-tuned model for ETA / RUL",
        "Fine-tuned models win on domain time-series; llm.py is model-agnostic"],
       ["Cloud", "Bedrock / Vertex managed inference",
        "Data residency and security; no self-run GPUs; OT data stays on-prem"],
       ["Platform", "Python with LangGraph / LangChain (or Helix no-code)",
        "Richest agent ecosystem; MVP is pure Python and runs offline"],
       ["Standards", "MCP for tools, A2A for agents, Langfuse for tracing",
        "Keeps integrations linear and every decision auditable"]],
      col_widths=[Inches(2.1), Inches(4.7), Inches(5.3)], fs=12, row_h=0.66)
footer_tag(s, "Architecture (10 pts) + Feasibility (15 pts)", "Handout 8 | docs/02 sec 8")
notes(s, "(75s) Straight from docs/02 section 8. The key message: strongest reasoning model, "
          "model-agnostic integration so we can swap, hosted in our own cloud for data "
          "control, Python on open standards so nothing is locked in.")

# =====================================================================
# SLIDE 11 - Benefits & next steps (Assignment)
# =====================================================================
s = slide(); sidebar(s, "BENEFITS & NEXT STEPS")
title_band(s, "11. Benefits, desired outcomes and next steps")
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.45),
    [[("Desired outcomes:", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.6), [
    "Protect the 14 million dollar quarterly loss and cut expedited freight by about a third.",
    "Turn days of analysis into minutes, with full Tier-2/3 visibility on critical parts.",
], size=15, gap=7)
txt(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.45),
    [[("Next steps:", 15, True, NAVY)]])
bullets(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.2), [
    ("Pilot: ", "a 3-month run on the Stuttgart line; measure stoppage hours and freight cost."),
    ("Learn: ", "feed planner approvals and edits back to tune thresholds and supplier rankings."),
    ("Scale: ", "extend to all 28 plants, then to the other case challenges (maintenance, quality) via the hierarchical growth path."),
    ("Harden: ", "live ERP/TMS/news connectors and full MCP tool servers, with Langfuse tracing in production."),
], size=14.5, gap=8)
rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.65), NAVY)
txt(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.65),
    [[("The ask: ", 15, True, GREEN),
      ("approve a one-line, one-quarter pilot. Low cost, low risk, clearly measurable. Thank you.",
       15, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer_tag(s, "Clarity, Presentation Quality & Creativity (20 pts)", "Assignment: benefits + next steps")
notes(s, "(60s) Close on outcomes then the roadmap. The hierarchical pattern from slide 4 "
          "is the growth path to the other case challenges. End on the pilot ask and take "
          "questions.")

# =====================================================================
out = "Supply_Chain_Sentinel_Design.pptx"
prs.save(out)

# --- verify: no em-dashes anywhere ---
from pptx import Presentation as _P
check = _P(out)
bad = 0
for sl in check.slides:
    for shp in sl.shapes:
        if shp.has_text_frame and ("—" in shp.text_frame.text or "–" in shp.text_frame.text):
            bad += 1
        if shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    if "—" in cell.text or "–" in cell.text:
                        bad += 1
    if sl.has_notes_slide and "—" in sl.notes_slide.notes_text_frame.text:
        bad += 1
print(f"Wrote {out} with {len(prs.slides._sldIdLst)} slides | em/en-dashes: {bad}")
