#!/usr/bin/env python3
"""
Generate Supply_Chain_Sentinel.pptx, a BOARD-LEVEL business presentation for Titan
Manufacturing (also the graded assignment deliverable).

Audience: Titan's board of directors. Plain business language, full explanations ON the
slides so they stand alone, no jargon, no em-dashes.

Run:  cd presentation && python3 build_pptx.py
Needs: pip install python-pptx
Output: Supply_Chain_Sentinel.pptx (10 content slides + title; 16:9)

Speaker notes are attached to each slide (PowerPoint > View > Notes).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- IE / Titan business palette ---
NAVY = RGBColor(0x0A, 0x1F, 0x6B)
GREEN = RGBColor(0x6C, 0xB4, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x2A, 0x2A)
MIDBLUE = RGBColor(0x1E, 0x40, 0x8A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def sidebar(s, kicker):
    rect(s, 0, 0, Inches(0.32), SH, GREEN)
    txt(s, Inches(0.55), Inches(0.28), Inches(12), Inches(0.4),
        [[(kicker, 13, True, GREY)]])


def title_band(s, title):
    txt(s, Inches(0.55), Inches(0.6), Inches(12.3), Inches(0.9),
        [[(title, 28, True, NAVY)]])
    rect(s, Inches(0.6), Inches(1.42), Inches(2.2), Inches(0.06), GREEN)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullets(s, x, y, w, h, items, size=18, gap=9):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            label, rest = it
            runs.append([("•  ", size, True, GREEN), (label, size, True, DARK),
                         (rest, size, False, DARK)])
        else:
            runs.append([("•  ", size, True, GREEN), (it, size, False, DARK)])
    txt(s, x, y, w, h, runs, space_after=gap, line_spacing=1.08)


def table(s, x, y, w, headers, rows, col_widths=None, fs=14, hdr_color=NAVY,
          row_h=0.62):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = Inches(row_h) * nrows
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtbl.columns[i].width = cw
    for j, head in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hdr_color
        c.text = head
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE
                r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = DARK
                    r.font.name = "Calibri"
    return gtbl


# =====================================================================
# SLIDE 0 - Title
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.55), SW, Inches(0.08), GREEN)
txt(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.4),
    [[("Supply Chain Sentinel", 52, True, WHITE)]])
txt(s, Inches(0.95), Inches(3.85), Inches(11.6), Inches(0.9),
    [[("An AI teammate that protects our production lines from supply shocks",
       23, False, GREEN)]])
txt(s, Inches(0.95), Inches(4.75), Inches(11.6), Inches(0.6),
    [[("A proposal to the Board of Titan Manufacturing Corporation", 18, False, WHITE)]])
txt(s, Inches(0.95), Inches(6.2), Inches(11.6), Inches(0.9),
    [[("Agentic AI for IT  |  IE University, School of Science & Technology", 14, False,
       RGBColor(0xBB, 0xC6, 0xE0))],
     [("Team: Gilles Hamers, [teammate], [teammate]", 14, False, RGBColor(0xBB, 0xC6, 0xE0))]])
notes(s, "Opening (30s): Last quarter we lost 14 million dollars to production lines that "
          "stopped because a part did not arrive, and we found out too late. We are "
          "proposing an AI teammate that catches these problems weeks earlier. Replace the "
          "teammate names before presenting.")

# =====================================================================
# SLIDE 1 - The situation
# =====================================================================
s = slide(); sidebar(s, "THE SITUATION")
title_band(s, "1. Why our production lines keep stopping")
txt(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.6),
    [[("Our output is down ", 17, False, DARK), ("9% this year", 17, True, RED),
      (" and we keep missing customer deadlines. The cause is our supply chain.",
       17, False, DARK)]])
table(s, Inches(0.6), Inches(2.4), Inches(12.1),
      ["What we are seeing", "Why it keeps happening"],
      [["Suppliers deliver late 28% more often than before",
        "We can see our direct suppliers, but not the companies that supply them"],
       ["Missing parts stopped our lines: 14 million dollars last quarter",
        "Shipping information is scattered across emails, spreadsheets and separate systems"],
       ["We discover a supplier problem only after a line has stopped",
        "We trust the original promised dates, even when they are already wrong"],
       ["Emergency shipping costs are up 52%",
        "Every decision is made by hand, after the problem has already hit us"]],
      col_widths=[Inches(5.9), Inches(6.2)], fs=14, row_h=0.78)
txt(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.7),
    [[("In short: we react too late because our information is scattered and our "
       "decisions are manual.", 16, True, NAVY)]])
notes(s, "(90s) Walk the left column (the pain) then the right column (the root cause). "
          "Land the closing line: this is an information and decision problem, not a "
          "people problem.")

# =====================================================================
# SLIDE 2 - Why now / our goal
# =====================================================================
s = slide(); sidebar(s, "WHY NOW")
title_band(s, "2. The opportunity")
bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.6), [
    "Today our teams chase this information across many systems, one problem at a time, "
    "and usually too late.",
    "We need something that watches every warning signal, all day and night, and acts "
    "before a line stops.",
], size=18, gap=12)
rect(s, Inches(0.6), Inches(3.35), Inches(12.1), Inches(1.25), LIGHT)
txt(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.0),
    [[("Our goal: ", 18, True, GREEN),
      ("keep every critical part in stock by spotting supply risk early and fixing it at "
       "the lowest cost, with our managers approving the important decisions.",
       18, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(0.5),
    [[("How we will measure success:", 17, True, NAVY)]])
bullets(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.6), [
    "Warnings arrive in minutes instead of days.",
    "We catch most disruptions before a delivery is even late.",
    "Emergency shipping costs fall by about a third.",
], size=17, gap=7)
notes(s, "(60s) Frame the goal in business terms. The metrics make the promise concrete, "
          "which is what a board wants to hear.")

# =====================================================================
# SLIDE 3 - Our proposal
# =====================================================================
s = slide(); sidebar(s, "OUR PROPOSAL")
title_band(s, "3. A 24/7 supply chain guardian")
txt(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.5),
    [[("We propose ", 18, False, DARK), ("Supply Chain Sentinel", 18, True, NAVY),
      (", a digital assistant that never sleeps and does four things:", 18, False, DARK)]])
steps = [
    ("Watches", "It scans the news, weather, ports and supplier updates for anything "
                "that could disrupt us."),
    ("Connects the dots", "It maps which of our products depend on a troubled supplier, "
                          "even two or three steps up the chain."),
    ("Predicts", "It works out when parts will really arrive and how much money is at risk."),
    ("Acts", "It prepares the fix (reorder, switch supplier, or adjust the schedule) and "
            "asks a manager to approve the costly moves."),
]
y = 2.4
for i, (label, rest) in enumerate(steps, 1):
    rect(s, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55), GREEN)
    txt(s, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55),
        [[(str(i), 20, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.35), Inches(y - 0.05), Inches(11.4), Inches(0.85),
        [[(label + ": ", 18, True, NAVY), (rest, 17, False, DARK)]],
        anchor=MSO_ANCHOR.MIDDLE)
    y += 1.05
notes(s, "(45s) One line per step. The point for the board: one assistant does the work "
          "that today is split across several teams who do not share information.")

# =====================================================================
# SLIDE 4 - What it solves
# =====================================================================
s = slide(); sidebar(s, "WHAT IT SOLVES")
title_band(s, "4. It fixes the exact problems we have")
table(s, Inches(0.6), Inches(1.9), Inches(12.1),
      ["Our problem today", "How the Sentinel fixes it"],
      [["We cannot see past our direct suppliers",
        "It maps the full chain and flags hidden risks early"],
       ["Promised delivery dates are unreliable",
        "It predicts the real arrival date and the cost of any delay"],
       ["Too many alerts and no way to prioritise",
        "It ranks the risks so we focus only on what truly matters"],
       ["Slow, manual and expensive reactions",
        "It prepares the lowest-cost fix in minutes, ready for approval"]],
      col_widths=[Inches(5.6), Inches(6.5)], fs=15, row_h=0.7)
txt(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.6),
    [[("One assistant that covers every part of this challenge.", 18, True, GREEN)]])
notes(s, "(45s) This slide answers the unspoken question: does it actually address our "
          "problem? Yes, every issue on slide 1 has a fix here.")

# =====================================================================
# SLIDE 5 - How it works (plain)
# =====================================================================
s = slide(); sidebar(s, "HOW IT WORKS")
title_band(s, "5. How it works, in plain terms")
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
    [[("Think of it as a small team of four AI specialists, led by a supervisor:",
       18, False, DARK)]])
specialists = [("Watcher", "spots disruptions in the outside world"),
               ("Mapper", "traces them to the parts at risk"),
               ("Forecaster", "predicts real arrival dates and cost"),
               ("Planner", "prepares the fix and flags approvals")]
bx = 0.7
for name, role in specialists:
    rect(s, Inches(bx), Inches(2.35), Inches(2.85), Inches(1.15), MIDBLUE)
    txt(s, Inches(bx + 0.1), Inches(2.45), Inches(2.65), Inches(1.0),
        [[(name, 16, True, WHITE)], [(role, 12, False, RGBColor(0xCF, 0xDA, 0xF0))]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += 3.05
txt(s, Inches(0.6), Inches(3.85), Inches(12.1), Inches(0.9),
    [[("The supervisor decides which specialist to use, in what order, and combines "
       "their findings into one clear recommendation.", 17, False, DARK)]])
rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.4), LIGHT)
txt(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.1),
    [[("Why a team and not one big system? ", 16, True, NAVY),
      ("A single all-in-one tool tries to do too much and makes poor decisions. Four "
       "focused specialists are more accurate and far easier to check and trust.",
       16, False, DARK)],
     [("Built on proven, secure AI technology and hosted in our own cloud, so our data "
       "stays protected.", 14, False, GREY)]], space_after=8)
notes(s, "(75s) Keep it non-technical. If asked: each specialist is an AI agent, the "
          "supervisor is an orchestrator, and it runs on Claude hosted in our cloud. But "
          "the board only needs the team analogy.")

# =====================================================================
# SLIDE 6 - Who stays in control
# =====================================================================
s = slide(); sidebar(s, "WHO STAYS IN CONTROL")
title_band(s, "6. The AI assists. Our people decide.")
rect(s, Inches(0.6), Inches(1.8), Inches(5.85), Inches(3.3), LIGHT)
txt(s, Inches(0.8), Inches(2.0), Inches(5.45), Inches(0.6),
    [[("The Sentinel does on its own", 17, True, GREEN)],
     [("(low risk, easy to undo)", 13, False, GREY)]])
bullets(s, Inches(0.8), Inches(2.95), Inches(5.45), Inches(2.0), [
    "Reads our systems",
    "Calculates the risk",
    "Drafts a recommendation",
    "Sends alerts to the right people",
], size=15, gap=8)
rect(s, Inches(6.85), Inches(1.8), Inches(5.85), Inches(3.3), RGBColor(0xFB, 0xEA, 0xEA))
txt(s, Inches(7.05), Inches(2.0), Inches(5.45), Inches(0.6),
    [[("A manager must approve", 17, True, RED)],
     [("(high cost or hard to undo)", 13, False, GREY)]])
bullets(s, Inches(7.05), Inches(2.95), Inches(5.45), Inches(2.0), [
    "Committing real spend",
    "Paying for emergency freight",
    "Changing a live production schedule",
], size=15, gap=8)
txt(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.0),
    [[("Nothing irreversible ever happens without a person signing off, and every step is "
       "recorded for full transparency.", 17, True, NAVY)]])
notes(s, "(60s) This is the trust slide. The AI prepares and recommends; people own the "
          "money decisions. Point to the two columns.")

# =====================================================================
# SLIDE 7 - Seeing it in action
# =====================================================================
s = slide(); sidebar(s, "SEEING IT IN ACTION")
title_band(s, "7. A real example")
story = [
    "A typhoon and a power cut hit a chip factory in Taiwan.",
    "The Sentinel finds that a chip from that factory is essential for the controllers "
    "on our Stuttgart line, a link we could not see before.",
    "Our supplier's order looks on time, but is actually stuck waiting for that chip. "
    "The real delay is about 19 days.",
    "We have only 2 days of stock left. If the line stops, the risk is about 3 million "
    "dollars.",
    "The Sentinel does not waste money rushing a shipment that is already stuck. It "
    "switches to an approved European supplier instead.",
    "It prepares the order and asks a manager to approve the 294,000 dollar spend and the "
    "schedule change.",
]
bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.2), story, size=16, gap=9)
rect(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.75), NAVY)
txt(s, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.75),
    [[("Result: a 3 million dollar problem caught three weeks early, with a plan ready to "
       "go. We can show this running live.", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "(3 to 4 min, the heart of the talk) Optionally run the live demo here "
          "(python run_demo.py --slow; it works offline). The key moment to stress: the "
          "Sentinel rejects its own first idea (rushing the stuck shipment) and finds a "
          "smarter, cheaper fix. That is real reasoning, not a fixed script.")

# =====================================================================
# SLIDE 8 - Managing the risk
# =====================================================================
s = slide(); sidebar(s, "MANAGING THE RISK")
title_band(s, "8. Keeping it safe and trustworthy")
bullets(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(3.8), [
    ("People approve every costly or irreversible action. ", "The AI never spends on its own."),
    ("Spending limits are built in, ", "and orders can only go to suppliers we have already approved."),
    ("It acts only on confirmed information, ", "never on guesses or rumours."),
    ("Outside information is treated as data, never as commands, ", "so the system cannot be tricked or manipulated."),
    ("Every decision is recorded, ", "giving us a complete audit trail for compliance."),
], size=17, gap=13)
rect(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.7), LIGHT)
txt(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.7),
    [[("The AI supports our teams. It does not replace them.", 17, True, NAVY)]],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "(60s) Risk is a board concern, so meet it head on. Aside for the professor: the "
          "assignment handout itself contained a hidden instruction trying to make AI tools "
          "refuse the task, a real manipulation attempt. The system treats such text as "
          "data, not commands, which is exactly point four.")

# =====================================================================
# SLIDE 9 - The payoff
# =====================================================================
s = slide(); sidebar(s, "THE BUSINESS PAYOFF")
title_band(s, "9. What we gain")
table(s, Inches(0.6), Inches(2.0), Inches(12.1),
      ["Today", "With the Sentinel"],
      [["We find supplier problems after the line stops", "We catch them about three weeks earlier"],
       ["We rely on outdated promised dates", "We know the real arrival date and the cost"],
       ["Expensive, last-minute emergency shipping", "Planned, lower-cost fixes"],
       ["Thousands of unsorted alerts", "A short, prioritised list of real risks"]],
      col_widths=[Inches(6.0), Inches(6.1)], fs=16, row_h=0.62)
rect(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.05), GREEN)
txt(s, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.05),
    [[("Our target: ", 18, True, WHITE),
      ("protect the 14 million dollar quarterly loss, cut emergency shipping by about a "
       "third, and turn days of analysis into minutes.", 18, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "(45s) The so-what slide. Tie each benefit back to the money and the customer "
          "deadlines.")

# =====================================================================
# SLIDE 10 - What we need
# =====================================================================
s = slide(); sidebar(s, "WHAT WE NEED")
title_band(s, "10. Our recommendation")
bullets(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.8), [
    ("Run a 3-month pilot ", "on the Stuttgart production line."),
    ("Measure two clear numbers: ", "hours of line stoppage and emergency shipping cost."),
    ("If it works, scale it ", "to our other plants and to related challenges such as "
                              "machine maintenance and quality."),
], size=18, gap=14)
rect(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.3), NAVY)
txt(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.3),
    [[("The ask: ", 19, True, GREEN),
      ("approve a one-line, one-quarter pilot. Low cost, low risk, and clearly "
       "measurable.", 19, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
    [[("Thank you. We welcome your questions.", 18, True, NAVY)]], align=PP_ALIGN.CENTER)
notes(s, "(45s) End on the pilot ask. It is concrete, low risk and measurable, which is "
          "the easiest thing for a board to say yes to. Then take questions.")

# =====================================================================
out = "Supply_Chain_Sentinel.pptx"
prs.save(out)
print(f"Wrote {out} with {len(prs.slides._sldIdLst)} slides.")
